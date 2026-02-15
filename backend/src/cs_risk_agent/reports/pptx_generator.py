"""PPTX レポート生成.

python-pptxを使用してリスク分析レポートをPowerPoint形式で生成する。
"""

from __future__ import annotations

import io
import logging
from datetime import datetime
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)

# リスクレベル色
_RISK_COLORS = {
    "critical": RGBColor(0xDC, 0x26, 0x26),
    "high": RGBColor(0xF9, 0x73, 0x16),
    "medium": RGBColor(0xEA, 0xB3, 0x08),
    "low": RGBColor(0x22, 0xC5, 0x5E),
}


def generate_risk_report_pptx(
    companies: list[dict[str, Any]],
    risk_scores: list[dict[str, Any]],
    alerts: list[dict[str, Any]],
    summary: dict[str, Any],
    fiscal_year: int,
    language: str = "ja",
) -> bytes:
    """リスク分析レポートPPTXを生成.

    Args:
        companies: 企業一覧
        risk_scores: リスクスコア一覧
        alerts: アラート一覧
        summary: リスクサマリー
        fiscal_year: 対象会計年度
        language: 言語 (ja/en)

    Returns:
        PPTX バイナリデータ
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- スライド1: タイトル ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_text_box(
        slide, 1, 2.5, 11.3, 1.5,
        "Consolidated Subsidiary Risk Analysis Report",
        font_size=28, bold=True, color=RGBColor(0x1E, 0x29, 0x3B),
    )
    _add_text_box(
        slide, 1, 4.0, 11.3, 0.8,
        f"Fiscal Year {fiscal_year} | Generated: {datetime.now().strftime('%Y-%m-%d')}",
        font_size=14, color=RGBColor(0x64, 0x74, 0x8B),
    )

    # --- スライド2: エグゼクティブサマリー ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Executive Summary")

    by_level = summary.get("by_level", {})
    total = summary.get("total_companies", 0)
    avg = summary.get("avg_score", 0)

    kpis = [
        ("Total Subsidiaries", str(total), RGBColor(0x3B, 0x82, 0xF6)),
        ("Average Risk Score", f"{avg:.1f}", RGBColor(0xF9, 0x73, 0x16)),
        ("Critical", str(by_level.get("critical", 0)), _RISK_COLORS["critical"]),
        ("High Risk", str(by_level.get("high", 0)), _RISK_COLORS["high"]),
    ]
    for i, (label, value, color) in enumerate(kpis):
        left = 0.8 + i * 3.1
        _add_kpi_box(slide, left, 2.0, 2.8, 1.6, label, value, color)

    # --- スライド3: リスクスコア一覧 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Risk Scores by Subsidiary")

    if risk_scores:
        sorted_scores = sorted(
            risk_scores, key=lambda x: x.get("total_score", 0), reverse=True,
        )
        rows = min(len(sorted_scores), 15)
        cols = 6
        table_shape = slide.shapes.add_table(
            rows + 1, cols,
            Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.4 * (rows + 1)),
        )
        table = table_shape.table

        headers = ["Entity Name", "Total Score", "DA", "Fraud", "Rule", "Level"]
        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = h
            _style_cell(
                cell, bold=True,
                bg=RGBColor(0x1E, 0x29, 0x3B),
                fg=RGBColor(0xFF, 0xFF, 0xFF),
            )

        for i, rs in enumerate(sorted_scores[:rows]):
            name = rs.get("entity_name", "")
            if len(name) > 30:
                name = name[:27] + "..."
            level = rs.get("risk_level", "low")
            data = [
                name,
                f"{rs.get('total_score', 0):.1f}",
                f"{rs.get('da_score', 0):.1f}",
                f"{rs.get('fraud_score', 0):.1f}",
                f"{rs.get('rule_score', 0):.1f}",
                level.upper(),
            ]
            for j, val in enumerate(data):
                cell = table.cell(i + 1, j)
                cell.text = val
                fg = _RISK_COLORS.get(level, RGBColor(0, 0, 0)) if j == 5 else None
                _style_cell(cell, font_size=10, fg=fg)

    # --- スライド4: アラート ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Risk Alerts")

    y_pos = 1.8
    for alert in alerts[:8]:
        sev = alert.get("severity", "medium")
        color = _RISK_COLORS.get(sev, RGBColor(0x64, 0x74, 0x8B))
        title = f"[{sev.upper()}] {alert.get('title', '')}"
        if len(title) > 80:
            title = title[:77] + "..."
        _add_text_box(slide, 0.8, y_pos, 11.7, 0.4, title, font_size=11, bold=True, color=color)
        y_pos += 0.5
        desc = alert.get("description", "")
        if len(desc) > 120:
            desc = desc[:117] + "..."
        _add_text_box(
            slide, 1.0, y_pos, 11.5, 0.3, desc,
            font_size=9, color=RGBColor(0x64, 0x74, 0x8B),
        )
        y_pos += 0.5

    # 出力
    buffer = io.BytesIO()
    prs.save(buffer)
    pptx_bytes = buffer.getvalue()
    buffer.close()
    logger.info("PPTX report generated: %d bytes", len(pptx_bytes))
    return pptx_bytes


def _add_text_box(
    slide: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    font_size: int = 12,
    bold: bool = False,
    color: RGBColor | None = None,
) -> None:
    """テキストボックスを追加."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color


def _add_title(slide: Any, text: str) -> None:
    """スライドタイトルを追加."""
    _add_text_box(
        slide, 0.5, 0.3, 12.3, 0.8, text,
        font_size=22, bold=True, color=RGBColor(0x1E, 0x29, 0x3B),
    )


def _add_kpi_box(
    slide: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    label: str,
    value: str,
    color: RGBColor,
) -> None:
    """KPIカード風ボックスを追加."""
    shape = slide.shapes.add_shape(
        1,  # rectangle
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = color

    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
    p2.alignment = PP_ALIGN.CENTER


def _style_cell(
    cell: Any,
    font_size: int = 9,
    bold: bool = False,
    bg: RGBColor | None = None,
    fg: RGBColor | None = None,
) -> None:
    """テーブルセルのスタイル設定."""
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        if fg:
            paragraph.font.color.rgb = fg
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
